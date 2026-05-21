#!/usr/bin/env -S uv run
# /// script
# requires-python = ">=3.11"
# dependencies = ["requests", "pdf2image", "python-pptx", "pyyaml"]
# ///
"""
Export slides to PPTX via the dev server API.
Requires: npm run dev (http://localhost:5173) to be running.

Usage:
  uv run scripts/export_pptx.py <project> [--port 5173] [--dpi 150]
  npm run export -- <project>
"""

import sys
import os
import json
import argparse
from pathlib import Path

ROOT = Path(__file__).parent.parent
OUTPUT_DIR = ROOT / "output"


def fetch_pdf(project: str, port: int) -> bytes:
    import requests
    url = f"http://localhost:{port}/api/projects/{project}/pdf"
    print(f"Requesting PDF from {url} ...")
    resp = requests.post(url, timeout=120)
    if resp.status_code != 200:
        print(f"Error: server returned {resp.status_code}: {resp.text}")
        sys.exit(1)
    print(f"  PDF received ({len(resp.content):,} bytes)")
    return resp.content


def load_notes(project: str) -> list[str]:
    """Return per-slide note texts in slide order. Empty string if no note."""
    project_dir = ROOT / "projects" / project

    slides_path = project_dir / "slides.json"
    if not slides_path.exists():
        return []
    with open(slides_path, encoding="utf-8") as f:
        slides = json.load(f)
    slide_ids = [s.get("id", "") for s in slides]

    script_path = project_dir / "script.yaml"
    if not script_path.exists():
        return [""] * len(slide_ids)

    import yaml
    with open(script_path, encoding="utf-8") as f:
        script: dict = yaml.safe_load(f) or {}

    notes = [script.get(sid, "").strip() for sid in slide_ids]
    filled = sum(1 for n in notes if n)
    print(f"  Script notes: {filled}/{len(notes)} slides have notes")
    return notes


def pdf_to_pptx(pdf_bytes: bytes, output_path: Path, dpi: int, notes: list[str]):
    from pdf2image import convert_from_bytes
    from pptx import Presentation
    from pptx.util import Inches
    import tempfile

    print(f"Converting PDF → PPTX (dpi={dpi}) ...")
    pages = convert_from_bytes(pdf_bytes, dpi=dpi)
    print(f"  {len(pages)} pages")

    w_px, h_px = pages[0].size
    slide_w = Inches(10)
    slide_h = int(slide_w * (h_px / w_px))

    prs = Presentation()
    prs.slide_width = slide_w
    prs.slide_height = slide_h
    blank_layout = prs.slide_layouts[6]

    with tempfile.TemporaryDirectory() as tmpdir:
        for i, page in enumerate(pages):
            img_path = os.path.join(tmpdir, f"page_{i:04d}.png")
            page.save(img_path, "PNG")
            slide = prs.slides.add_slide(blank_layout)
            slide.shapes.add_picture(img_path, 0, 0, slide_w, slide_h)
            if i < len(notes) and notes[i]:
                slide.notes_slide.notes_text_frame.text = notes[i]
            print(f"  [{i+1}/{len(pages)}] added", end="\r")

    prs.save(str(output_path))
    print(f"\nSaved: {output_path}")


def main():
    parser = argparse.ArgumentParser(description=__doc__, formatter_class=argparse.RawDescriptionHelpFormatter)
    parser.add_argument("project", help="Project name (directory under projects/)")
    parser.add_argument("--port", type=int, default=5173, help="Dev server port (default: 5173)")
    parser.add_argument("--dpi", type=int, default=150, help="Render DPI (default: 150, higher = better quality)")
    args = parser.parse_args()

    OUTPUT_DIR.mkdir(exist_ok=True)

    pdf_bytes = fetch_pdf(args.project, args.port)

    pdf_path = OUTPUT_DIR / f"{args.project}.pdf"
    pdf_path.write_bytes(pdf_bytes)
    print(f"  PDF saved: {pdf_path}")

    notes = load_notes(args.project)

    pptx_path = OUTPUT_DIR / f"{args.project}.pptx"
    pdf_to_pptx(pdf_bytes, pptx_path, args.dpi, notes)


if __name__ == "__main__":
    main()
